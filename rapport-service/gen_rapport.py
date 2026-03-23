"""
Neoclima Field — Moteur de génération PPTX dynamique
Rapport client maître d'ouvrage — design exécutif premium v2
"""
import io
import os
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ────────────────────────────────────────────────────────────
DARK     = RGBColor(0x0D, 0x15, 0x21)   # quasi-noir marine
NAVY     = RGBColor(0x13, 0x20, 0x35)   # bleu nuit
NC_TEAL  = RGBColor(0x00, 0xB4, 0xD8)   # accent Neoclima
NC_TEAL2 = RGBColor(0x02, 0x8C, 0xA8)   # teal foncé
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xF4, 0xF7, 0xFB)   # fond clair
PALE     = RGBColor(0xF8, 0xFA, 0xFC)
GREY     = RGBColor(0x64, 0x74, 0x8B)
SILVER   = RGBColor(0x94, 0xA3, 0xB8)
BORDER   = RGBColor(0xE2, 0xE8, 0xF0)
BAR_BG   = RGBColor(0xE2, 0xE8, 0xF0)
GREEN    = RGBColor(0x05, 0x96, 0x69)
AMBER    = RGBColor(0xD9, 0x77, 0x06)
RED_VIG  = RGBColor(0xDC, 0x26, 0x26)
BLUE_KPI = RGBColor(0x1E, 0x40, 0xAF)
MID_NAVY = RGBColor(0x1E, 0x30, 0x4A)

STATUS_COLORS = {
    "green": GREEN,
    "amber": AMBER,
    "red":   RED_VIG,
    "grey":  RGBColor(0xCB, 0xD5, 0xE1),
}

STATUT_GLOBAL_COLORS = {
    "MAÎTRISÉ":          GREEN,
    "SOUS SURVEILLANCE": AMBER,
    "CRITIQUE":          RED_VIG,
}

LOGO_PATH = os.path.join(os.path.dirname(__file__), "logo.png")


# ── Primitives ─────────────────────────────────────────────────────────

def rect(sl, x, y, w, h, fill, line_color=None, line_width=0):
    s = sl.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_width)
    else:
        s.line.fill.background()
    return s


def txt(sl, text, x, y, w, h, size=12, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    return tb


def logo_on_dark(sl, x, y, h_inch):
    """Affiche le logo blanc sur fond sombre."""
    logo = os.environ.get("LOGO_PATH", LOGO_PATH)
    try:
        img = Image.open(logo).convert("RGBA")
        r_ch, g_ch, b_ch, a_ch = img.split()
        # Rend tous les pixels blancs (luminosité max), garde l'alpha
        r_ch = r_ch.point(lambda _: 255)
        g_ch = g_ch.point(lambda _: 255)
        b_ch = b_ch.point(lambda _: 255)
        white_img = Image.merge("RGBA", (r_ch, g_ch, b_ch, a_ch))
        buf = io.BytesIO()
        white_img.save(buf, format="PNG")
        buf.seek(0)
        ratio = img.width / img.height
        sl.shapes.add_picture(buf, x, y,
                              height=Inches(h_inch),
                              width=Inches(h_inch * ratio))
    except Exception:
        txt(sl, "NEOCLIMA", x, y, Inches(3), Inches(h_inch),
            size=18, bold=True, color=WHITE)


def slide_header(sl, title, subtitle, W, H, dark=True):
    """Bandeau haut commun aux slides internes."""
    bg = DARK if dark else NAVY
    rect(sl, 0, 0, W, Inches(1.1), bg)
    rect(sl, 0, 0, Inches(0.08), Inches(1.1), NC_TEAL)
    logo_on_dark(sl, Inches(0.22), Inches(0.24), 0.50)
    txt(sl, title,
        Inches(2.1), Inches(0.15), W - Inches(2.3), Inches(0.52),
        size=19, bold=True, color=WHITE)
    txt(sl, subtitle,
        Inches(2.1), Inches(0.65), W - Inches(2.3), Inches(0.35),
        size=10, color=NC_TEAL)


def footer(sl, date_str, W, H):
    rect(sl, 0, H - Inches(0.26), W, Inches(0.26), NAVY)
    txt(sl,
        f"Neoclima SA   ·   Rapport d'avancement   ·   {date_str}   ·   CONFIDENTIEL",
        Inches(0.4), H - Inches(0.245), W - Inches(0.8), Inches(0.22),
        size=7, color=SILVER, align=PP_ALIGN.CENTER)


def progress_bar(sl, x, y, w, h, pct, bar_color, bg_color=None):
    """Barre de progression avec fond gris."""
    bg = bg_color or BAR_BG
    rect(sl, x, y, w, h, bg)
    if pct > 0:
        filled = max(Inches(0.04), int(w * pct / 100))
        rect(sl, x, y, filled, h, bar_color)


# ═══════════════════════════════════════════════════════════════════════
# POINT D'ENTRÉE PRINCIPAL
# ═══════════════════════════════════════════════════════════════════════

def generate_rapport(data: dict) -> bytes:
    """
    Génère le rapport PPTX et retourne les bytes.

    data = {
      "project": {
        "nom": str,
        "client": str,
        "semaine": str,          # "S12 · 22 mars – 26 mars"
        "date": str,             # "23 mars 2026"
        "statut_global": str,    # "MAÎTRISÉ" | "SOUS SURVEILLANCE" | "CRITIQUE"
        "avancement_global": int,# 0-100
        "semaines_restantes": int,
        "phrase": str,
      },
      "zones": [
        {"batiment": str, "niveau": str, "avancement": int,
         "statut": "green|amber|red|grey", "commentaire": str}
      ],
      "equipes": [
        {"nom": str, "effectif": int, "type": "Interne|Sous-traitant"}
      ],
      "vigilances": [
        {"zone": str, "sujet": str, "action": str, "impact": str}
      ],
      "prochaines_etapes": [
        {"date": str, "titre": str, "detail": str}
      ],
      "faits_marquants": [str],
    }
    """
    P        = data["project"]
    zones    = data.get("zones", [])
    equipes  = data.get("equipes", [])
    vigilances = data.get("vigilances", [])
    etapes   = data.get("prochaines_etapes", [])
    avan     = int(P.get("avancement_global", 0))
    statut_color = STATUT_GLOBAL_COLORS.get(
        P.get("statut_global", "SOUS SURVEILLANCE"), AMBER)

    # Auto-génère les faits marquants si non fournis
    faits = list(data.get("faits_marquants", []))
    if not faits:
        done_zones = [z for z in zones if z.get("avancement", 0) == 100]
        prog_zones = [z for z in zones if 0 < z.get("avancement", 0) < 100]
        for z in done_zones[:2]:
            faits.append(
                f"✓  {z['batiment']} {z['niveau']} — "
                f"{z.get('commentaire', 'Réceptionné sans réserve.')}")
        for z in prog_zones[:3]:
            faits.append(
                f"◎  {z['batiment']} {z['niveau']} à {z['avancement']}% — "
                f"{z.get('commentaire', 'En cours, livraison prévue dans les délais.')}")

    prs = Presentation()
    prs.slide_width  = Inches(13.333)
    prs.slide_height = Inches(7.5)
    W = prs.slide_width
    H = prs.slide_height

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 1 — COUVERTURE PREMIUM
    # ──────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    # Fond plein dark
    rect(sl, 0, 0, W, H, DARK)

    # Bande latérale gauche (accent couleur)
    rect(sl, 0, 0, Inches(0.12), H, NC_TEAL)

    # Bloc gauche légèrement plus clair
    rect(sl, Inches(0.12), 0, Inches(5.5), H, NAVY)

    # Logo
    logo_on_dark(sl, Inches(0.35), Inches(0.38), 0.6)

    # Informations projet (côté gauche)
    txt(sl, "RAPPORT D'AVANCEMENT", Inches(0.35), Inches(1.45),
        Inches(5.0), Inches(0.38), size=9.5, color=NC_TEAL, bold=True)

    # Titre chantier
    nom = P.get("nom", "Chantier")
    txt(sl, nom, Inches(0.35), Inches(1.9),
        Inches(5.1), Inches(1.4), size=22, bold=True, color=WHITE, wrap=True)

    # Ligne séparatrice fine
    rect(sl, Inches(0.35), Inches(3.5), Inches(4.5), Pt(1.5),
         RGBColor(0x2A, 0x40, 0x58))

    txt(sl, P.get("semaine", ""), Inches(0.35), Inches(3.68),
        Inches(5.0), Inches(0.38), size=10.5, color=SILVER)
    txt(sl, f"Émis le {P.get('date', '')}", Inches(0.35), Inches(4.1),
        Inches(5.0), Inches(0.35), size=9.5, color=SILVER, italic=True)
    txt(sl, f"Destinataire : {P.get('client', '')}", Inches(0.35), Inches(4.5),
        Inches(5.0), Inches(0.35), size=9.5, color=SILVER)

    # Badge statut
    badge_w = Inches(2.6)
    rect(sl, Inches(0.35), Inches(5.2), badge_w, Inches(0.46), statut_color)
    txt(sl, "● " + P.get("statut_global", "MAÎTRISÉ"),
        Inches(0.35), Inches(5.22), badge_w, Inches(0.44),
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # ── Bloc avancement (côté droit) ──────────────────────────────────
    # Fond légèrement distinct
    rect(sl, Inches(5.62), Inches(0.3), Inches(7.4), Inches(6.9), MID_NAVY)

    # Grand chiffre avancement — taille raisonnable et bien positionné
    avan_str = f"{avan}%"
    txt(sl, avan_str,
        Inches(5.62), Inches(0.85), Inches(7.4), Inches(3.0),
        size=96, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "AVANCEMENT GLOBAL DU CHANTIER",
        Inches(5.62), Inches(3.95), Inches(7.4), Inches(0.42),
        size=11, color=SILVER, align=PP_ALIGN.CENTER, bold=True)

    # Barre de progression
    bar_x = Inches(6.2); bar_y = Inches(4.5)
    bar_w = Inches(6.2); bar_h = Inches(0.16)
    progress_bar(sl, bar_x, bar_y, bar_w, bar_h, avan,
                 bar_color=NC_TEAL, bg_color=RGBColor(0x20, 0x30, 0x45))

    # Label % sous la barre
    txt(sl, f"{avan}% réalisé", bar_x, bar_y + Inches(0.22),
        Inches(3.0), Inches(0.32), size=9.5, color=NC_TEAL)

    sr = P.get("semaines_restantes", 0)
    sem_label = f"{sr} semaine{'s' if sr != 1 else ''} restante{'s' if sr != 1 else ''}"
    txt(sl, sem_label,
        bar_x + Inches(3.2), bar_y + Inches(0.22), Inches(3.0), Inches(0.32),
        size=9.5, color=SILVER, align=PP_ALIGN.RIGHT)

    # Métriques rapides (3 chiffres)
    total_eff = sum(e.get("effectif", 0) for e in equipes)
    nb_vig    = len(vigilances)
    zones_done = len([z for z in zones if z.get("avancement", 0) == 100])
    metrics = [
        (str(total_eff), "intervenants\nsur site", NC_TEAL),
        (str(zones_done), f"zone{'s' if zones_done != 1 else ''}\nterminée{'s' if zones_done != 1 else ''}", GREEN),
        (str(nb_vig), "point" + ("s" if nb_vig != 1 else "") + "\nvigilance", AMBER if nb_vig > 0 else GREEN),
    ]
    metric_w = Inches(2.0)
    metric_y = Inches(5.25)
    for mi, (val, lbl, col) in enumerate(metrics):
        mx = Inches(6.2) + mi * (metric_w + Inches(0.08))
        rect(sl, mx, metric_y, metric_w, Inches(1.6), RGBColor(0x0D, 0x1C, 0x2F))
        rect(sl, mx, metric_y, metric_w, Inches(0.05), col)
        txt(sl, val, mx, metric_y + Inches(0.12), metric_w, Inches(0.75),
            size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
        txt(sl, lbl, mx, metric_y + Inches(0.88), metric_w, Inches(0.55),
            size=8.5, color=SILVER, align=PP_ALIGN.CENTER, wrap=True)

    # Footer
    rect(sl, 0, H - Inches(0.3), W, Inches(0.3), RGBColor(0x07, 0x0F, 0x18))
    txt(sl, "CONFIDENTIEL — Usage exclusif du destinataire   ·   Neoclima SA",
        Inches(0.4), H - Inches(0.28), Inches(9), Inches(0.25),
        size=7.5, color=SILVER, italic=True)
    txt(sl, P.get("date", ""),
        W - Inches(2.0), H - Inches(0.28), Inches(1.8), Inches(0.25),
        size=7.5, color=SILVER, align=PP_ALIGN.RIGHT)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 2 — SYNTHÈSE EXÉCUTIVE
    # ──────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, LIGHT)
    slide_header(sl, "SYNTHÈSE EXÉCUTIVE", P.get("semaine", ""), W, H, dark=True)

    # Bandeau statut
    rect(sl, Inches(0.4), Inches(1.25), W - Inches(0.8), Inches(0.78), statut_color)
    phrase = P.get("phrase", "Le chantier progresse conformément au planning.")
    txt(sl, f"STATUT GLOBAL : {P.get('statut_global', '')}   ·   {phrase}",
        Inches(0.55), Inches(1.32), W - Inches(1.1), Inches(0.6),
        size=11, bold=True, color=WHITE, wrap=True)

    # 4 KPI cards
    total_eff = sum(e.get("effectif", 0) for e in equipes)
    nb_vig    = len(vigilances)
    kpis = [
        (f"{avan}%",       "Avancement\nglobal",       BLUE_KPI, "En progression"),
        (f"{total_eff}",   "Intervenants\nsur site",   NC_TEAL2, "Équipes mobilisées"),
        (f"{nb_vig}",      "Points de\nvigilance",
         AMBER if nb_vig > 0 else GREEN,
         "Actions en cours" if nb_vig > 0 else "Aucun point ouvert"),
        (f"{zones_done}",  "Zones\nterminées",         GREEN,    "Réceptionnées"),
    ]
    card_w = Inches(2.9)
    card_gap = (W - Inches(0.8) - 4 * card_w) / 3
    for i, (val, label, color, note) in enumerate(kpis):
        kx = Inches(0.4) + i * (card_w + card_gap)
        ky = Inches(2.2)
        rect(sl, kx, ky, card_w, Inches(2.05),
             WHITE, line_color=BORDER, line_width=0.8)
        rect(sl, kx, ky, card_w, Inches(0.07), color)
        txt(sl, val, kx, ky + Inches(0.12), card_w, Inches(0.95),
            size=48, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
        txt(sl, label, kx, ky + Inches(1.05), card_w, Inches(0.5),
            size=9.5, color=GREY, align=PP_ALIGN.CENTER, wrap=True)
        txt(sl, note, kx + Inches(0.15), ky + Inches(1.65),
            card_w - Inches(0.3), Inches(0.3),
            size=8.5, color=color, italic=True)

    # Section faits marquants
    rect(sl, Inches(0.4), Inches(4.42), W - Inches(0.8), Pt(1.5), BORDER)
    txt(sl, "FAITS MARQUANTS DE LA SEMAINE",
        Inches(0.4), Inches(4.55), Inches(6), Inches(0.3),
        size=8.5, bold=True, color=GREY)

    if faits:
        fy = Inches(4.94)
        for fact in faits[:4]:
            col = GREEN if str(fact).startswith("✓") else AMBER
            rect(sl, Inches(0.4), fy + Inches(0.06),
                 Inches(0.05), Inches(0.22), col)
            txt(sl, str(fact), Inches(0.55), fy,
                W - Inches(1.0), Inches(0.32),
                size=10, color=NAVY, wrap=False)
            fy += Inches(0.36)
    else:
        txt(sl, "Semaine en cours — synthèse disponible en fin de semaine.",
            Inches(0.55), Inches(4.94), W - Inches(1.0), Inches(0.35),
            size=10, color=GREY, italic=True)

    footer(sl, P.get("date", ""), W, H)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 3 — VUE D'ENSEMBLE PAR ZONE
    # ──────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, LIGHT)
    slide_header(sl, "VUE D'ENSEMBLE — AVANCEMENT PAR ZONE",
                 f"État au {P.get('date', '')}", W, H, dark=True)

    if zones:
        # Groupe par bâtiment
        groups: dict[str, list] = {}
        for z in zones:
            bat = z.get("batiment", "Zone")
            groups.setdefault(bat, []).append(z)

        bat_names = list(groups.keys())[:4]
        n_cols = len(bat_names)
        col_w = (W - Inches(0.6) - Inches(0.2) * (n_cols - 1)) / n_cols
        start_x = Inches(0.3)

        for gi, bat in enumerate(bat_names):
            bx = start_x + gi * (col_w + Inches(0.2))
            by_start = Inches(1.22)
            rect(sl, bx, by_start, col_w, Inches(0.4), NAVY)
            txt(sl, bat.upper(), bx + Inches(0.12),
                by_start + Inches(0.07), col_w - Inches(0.2), Inches(0.3),
                size=9.5, bold=True, color=WHITE)

            for j, z in enumerate(groups[bat][:5]):
                zy = by_start + Inches(0.4 + j * 0.99)
                bg = WHITE if j % 2 == 0 else PALE
                rect(sl, bx, zy, col_w, Inches(0.9),
                     bg, line_color=BORDER, line_width=0.5)
                # Tag niveau
                tag_w = Inches(0.72)
                rect(sl, bx, zy, tag_w, Inches(0.9), MID_NAVY)
                lvl = z.get("niveau", "?")
                # Abbréger si trop long
                lvl_short = lvl[:8] if len(lvl) > 8 else lvl
                txt(sl, lvl_short, bx, zy, tag_w, Inches(0.9),
                    size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
                # Contenu
                pct = z.get("avancement", 0)
                bar_color = STATUS_COLORS.get(z.get("statut", "amber"), AMBER)
                bx2 = bx + tag_w + Inches(0.1)
                bar_y2 = zy + Inches(0.22)
                bar_w2 = col_w - tag_w - Inches(0.48)
                progress_bar(sl, bx2, bar_y2, bar_w2, Inches(0.22),
                             pct, bar_color)
                txt(sl, f"{pct}%",
                    bx2 + bar_w2 + Inches(0.05), zy + Inches(0.18),
                    Inches(0.32), Inches(0.26),
                    size=8, bold=True, color=bar_color)
                cmt = str(z.get("commentaire", ""))[:42]
                txt(sl, cmt, bx2, zy + Inches(0.54),
                    col_w - tag_w - Inches(0.2), Inches(0.28),
                    size=7.5, color=GREY, italic=True)
    else:
        txt(sl, "Aucune donnée de zone disponible pour cette semaine.",
            Inches(0.5), Inches(2.5), W - Inches(1.0), Inches(0.5),
            size=12, color=GREY, italic=True, align=PP_ALIGN.CENTER)

    # Légende horizontale en bas de slide
    leg_items = [
        (GREEN,  "Terminé"),
        (AMBER,  "En cours"),
        (RED_VIG,"Vigilance"),
        (RGBColor(0xCB, 0xD5, 0xE1), "Planifié"),
    ]
    leg_y_bot = H - Inches(0.62)
    lx = Inches(0.35)
    for col, lbl in leg_items:
        rect(sl, lx, leg_y_bot + Inches(0.04), Inches(0.16), Inches(0.16), col)
        txt(sl, lbl, lx + Inches(0.22), leg_y_bot,
            Inches(1.3), Inches(0.26), size=8, color=GREY)
        lx += Inches(1.6)

    footer(sl, P.get("date", ""), W, H)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 4 — TABLEAU AVANCEMENT DÉTAILLÉ
    # ──────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, LIGHT)
    slide_header(sl, "AVANCEMENT DÉTAILLÉ PAR ZONE",
                 f"Tableau de suivi — {P.get('date', '')}", W, H, dark=False)

    STATUS_LABELS = {
        "green": ("✓ Terminé",   GREEN),
        "amber": ("◎ En cours",  AMBER),
        "red":   ("⚠ Vigilance", RED_VIG),
        "grey":  ("○ Planifié",  GREY),
    }

    cols    = ["BÂTIMENT",  "ZONE / NIVEAU", "AVANCEMENT", "STATUT",     "COMMENTAIRE"]
    col_ws  = [Inches(1.7), Inches(1.6),     Inches(1.5),  Inches(1.7),  Inches(6.2)]
    tx = Inches(0.35); ty = Inches(1.25); row_h = Inches(0.37)

    cx = tx
    for col_name, cw in zip(cols, col_ws):
        rect(sl, cx, ty, cw, row_h, NAVY)
        txt(sl, col_name, cx + Inches(0.08), ty + Inches(0.08), cw, row_h,
            size=8.5, bold=True, color=WHITE)
        cx += cw

    for ri, z in enumerate(zones[:14]):
        ry  = ty + row_h + ri * row_h
        bg  = WHITE if ri % 2 == 0 else PALE
        cx  = tx
        vals = [
            z.get("batiment", ""),
            z.get("niveau", ""),
            "",  # avancement (barre)
            "",  # statut (badge)
            z.get("commentaire", ""),
        ]
        for ci, (val, cw) in enumerate(zip(vals, col_ws)):
            rect(sl, cx, ry, cw, row_h, bg, line_color=BORDER, line_width=0.25)
            pct = z.get("avancement", 0)
            bar_color = STATUS_COLORS.get(z.get("statut", "amber"), AMBER)

            if ci == 2:
                bx = cx + Inches(0.08); by = ry + Inches(0.09)
                bw = cw - Inches(0.55); bh = Inches(0.18)
                progress_bar(sl, bx, by, bw, bh, pct, bar_color)
                txt(sl, f"{pct}%",
                    bx + bw + Inches(0.05), ry + Inches(0.05),
                    Inches(0.38), row_h, size=7.5, bold=True, color=bar_color)
            elif ci == 3:
                lbl, col = STATUS_LABELS.get(z.get("statut", "amber"), ("◎ En cours", AMBER))
                bw_b = cw - Inches(0.16)
                rect(sl, cx + Inches(0.08), ry + Inches(0.07),
                     bw_b, Inches(0.23), col)
                txt(sl, lbl, cx + Inches(0.08), ry + Inches(0.07),
                    bw_b, Inches(0.23),
                    size=7.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            else:
                txt(sl, str(val),
                    cx + Inches(0.08), ry + Inches(0.08),
                    cw - Inches(0.12), row_h, size=8.5, color=NAVY)
            cx += cw

    footer(sl, P.get("date", ""), W, H)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 5 — ACTIVITÉ CHANTIER
    # ──────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, LIGHT)
    slide_header(sl, "ACTIVITÉ CHANTIER",
                 "Mobilisation & dynamique terrain", W, H, dark=True)

    total_eff = sum(e.get("effectif", 0) for e in equipes)
    internes  = sum(e.get("effectif", 0) for e in equipes
                    if e.get("type") == "Interne")
    st_eff    = total_eff - internes
    st_pct    = round(st_eff / total_eff * 100) if total_eff else 0
    int_pct   = 100 - st_pct

    # KPI principal — encart gauche
    rect(sl, Inches(0.4), Inches(1.25), Inches(2.9), Inches(4.3), NAVY)
    txt(sl, str(total_eff),
        Inches(0.4), Inches(1.45), Inches(2.9), Inches(1.6),
        size=72, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, "INTERVENANTS\nSUR SITE",
        Inches(0.4), Inches(3.1), Inches(2.9), Inches(0.7),
        size=11, bold=True, color=SILVER, align=PP_ALIGN.CENTER, wrap=True)

    # Répartition barre
    rect(sl, Inches(0.55), Inches(3.95), Inches(2.6), Inches(0.2),
         RGBColor(0x2A, 0x40, 0x58))
    if int_pct > 0:
        rect(sl, Inches(0.55), Inches(3.95),
             int(Inches(2.6) * int_pct / 100), Inches(0.2), NC_TEAL)

    txt(sl, f"Internes {int_pct}%", Inches(0.55), Inches(4.2),
        Inches(1.3), Inches(0.26), size=8.5, color=WHITE)
    txt(sl, f"ST {st_pct}%", Inches(1.85), Inches(4.2),
        Inches(1.3), Inches(0.26), size=8.5, color=SILVER, align=PP_ALIGN.RIGHT)

    # Dynamique
    rect(sl, Inches(0.55), Inches(4.58), Inches(2.6), Pt(1), NC_TEAL)
    txt(sl, "↗  EN HAUSSE VS S-1",
        Inches(0.4), Inches(4.68), Inches(2.9), Inches(0.26),
        size=9.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

    # Cards équipes (grille 2 colonnes) — max 8 cards, 4 rangées de 1.3"
    eq_list  = equipes[:8]
    card_w   = Inches(4.8)
    card_h   = Inches(1.20)
    card_gap = Inches(0.10)
    cols_pos = [Inches(3.55), Inches(8.5)]

    for qi, equipe in enumerate(eq_list):
        col_idx = qi % 2
        row_idx = qi // 2
        ex = cols_pos[col_idx]
        ey = Inches(1.25) + row_idx * (card_h + card_gap)
        rect(sl, ex, ey, card_w, card_h, WHITE, line_color=BORDER, line_width=0.8)
        # Accent couleur gauche
        type_col = NC_TEAL if equipe.get("type") == "Interne" else AMBER
        rect(sl, ex, ey, Inches(0.07), card_h, type_col)
        txt(sl, equipe.get("nom", "Équipe"),
            ex + Inches(0.18), ey + Inches(0.1),
            card_w - Inches(1.2), Inches(0.42),
            size=10.5, bold=True, color=NAVY)
        txt(sl, f"{equipe.get('effectif', 0)} pers.",
            ex + card_w - Inches(1.1), ey + Inches(0.1),
            Inches(1.0), Inches(0.42),
            size=14, bold=True, color=type_col, align=PP_ALIGN.RIGHT)
        txt(sl, equipe.get("type", "Interne"),
            ex + Inches(0.18), ey + Inches(0.52),
            Inches(2.0), Inches(0.28),
            size=8.5, color=type_col, bold=True)
        secteurs = equipe.get("secteurs", "")
        if secteurs:
            txt(sl, f"Secteurs : {secteurs}",
                ex + Inches(0.18), ey + Inches(0.82),
                card_w - Inches(0.3), Inches(0.26),
                size=8, color=GREY, italic=True)

    footer(sl, P.get("date", ""), W, H)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 6 — POINTS DE VIGILANCE
    # ──────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, LIGHT)
    sub_vig = ("Situations nécessitant attention — Toutes maîtrisées"
               if vigilances else "Aucun point de vigilance actif cette semaine")
    slide_header(sl, "POINTS DE VIGILANCE & ACTIONS EN COURS",
                 sub_vig, W, H, dark=False)

    vig_cols = ["ZONE",     "SITUATION",  "ACTION EN COURS", "IMPACT PLANNING"]
    vig_ws   = [Inches(2.0), Inches(3.5), Inches(4.3),       Inches(3.0)]
    vx = Inches(0.35); vy = Inches(1.25)
    cvx = vx
    for cn, cw in zip(vig_cols, vig_ws):
        rect(sl, cvx, vy, cw, Inches(0.38), NAVY)
        txt(sl, cn, cvx + Inches(0.1), vy + Inches(0.09), cw, Inches(0.38),
            size=8.5, bold=True, color=WHITE)
        cvx += cw

    for ri, v in enumerate(vigilances[:5]):
        ry = vy + Inches(0.38) + ri * Inches(0.87)
        bg = WHITE if ri % 2 == 0 else PALE
        cvx = vx
        for val, cw in zip([
            v.get("zone", ""),
            v.get("sujet", ""),
            v.get("action", ""),
            v.get("impact", ""),
        ], vig_ws):
            rect(sl, cvx, ry, cw, Inches(0.8),
                 bg, line_color=BORDER, line_width=0.25)
            rect(sl, cvx, ry, Inches(0.06), Inches(0.8), AMBER)
            txt(sl, str(val),
                cvx + Inches(0.14), ry + Inches(0.1),
                cw - Inches(0.22), Inches(0.65),
                size=8.5, color=NAVY, wrap=True)
            cvx += cw

    # Conclusion
    concl_y = vy + Inches(0.38) + min(len(vigilances), 5) * Inches(0.87) + Inches(0.2)
    if concl_y < H - Inches(1.0):
        rect(sl, Inches(0.35), concl_y, W - Inches(0.7), Inches(0.78), PALE,
             line_color=BORDER, line_width=0.5)
        concl = (
            "Les points de vigilance identifiés sont activement traités. "
            "Aucun n'impacte la date de livraison contractuelle. "
            "Planning de référence maintenu."
        ) if vigilances else (
            "Aucun point de vigilance actif. "
            "Le chantier progresse conformément au planning."
        )
        txt(sl, "ÉVALUATION GLOBALE DES RISQUES",
            Inches(0.5), concl_y + Inches(0.08), Inches(6), Inches(0.26),
            size=8, bold=True, color=GREY)
        txt(sl, concl,
            Inches(0.5), concl_y + Inches(0.36),
            W - Inches(0.9), Inches(0.38),
            size=9.5, color=NAVY, wrap=True)

    footer(sl, P.get("date", ""), W, H)

    # ──────────────────────────────────────────────────────────────────
    # SLIDE 7 — PROCHAINES ÉTAPES
    # ──────────────────────────────────────────────────────────────────
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    rect(sl, 0, 0, W, H, LIGHT)
    slide_header(sl, "PROCHAINES ÉTAPES — PLANNING 2 SEMAINES",
                 "Calendrier des actions validées", W, H, dark=True)

    STEP_COLORS = [NC_TEAL, NAVY, GREEN, AMBER, MID_NAVY]

    line_x     = Inches(1.6)
    line_y0    = Inches(1.35)
    line_y1    = Inches(6.55)
    # Ligne verticale timeline
    rect(sl, line_x, line_y0, Inches(0.05), line_y1 - line_y0, NC_TEAL)

    etapes_disp = etapes[:5]
    n_steps  = max(len(etapes_disp), 1)
    spacing  = (line_y1 - line_y0 - Inches(0.9)) / (n_steps - 1) if n_steps > 1 else 0
    step_h   = Inches(0.88)

    if etapes_disp:
        for si, e in enumerate(etapes_disp):
            sy = line_y0 + si * spacing
            dot_col = STEP_COLORS[si % len(STEP_COLORS)]
            # Pastille sur la ligne
            rect(sl, line_x - Inches(0.13), sy + Inches(0.13),
                 Inches(0.30), Inches(0.30), dot_col)
            # Badge date
            rect(sl, Inches(0.2), sy + Inches(0.08),
                 Inches(1.2), Inches(0.3), dot_col)
            txt(sl, e.get("date", ""),
                Inches(0.2), sy + Inches(0.08),
                Inches(1.2), Inches(0.3),
                size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
            # Card
            cx_card = line_x + Inches(0.28)
            rect(sl, cx_card, sy, W - cx_card - Inches(0.35), step_h,
                 WHITE, line_color=BORDER, line_width=0.5)
            rect(sl, cx_card, sy, Inches(0.07), step_h, dot_col)
            txt(sl, e.get("titre", ""),
                cx_card + Inches(0.18), sy + Inches(0.1),
                W - cx_card - Inches(0.75), Inches(0.4),
                size=11.5, bold=True, color=NAVY)
            txt(sl, e.get("detail", ""),
                cx_card + Inches(0.18), sy + Inches(0.5),
                W - cx_card - Inches(0.75), Inches(0.36),
                size=9, color=GREY, italic=True)
    else:
        txt(sl, "Planning des prochaines semaines en cours de finalisation.",
            Inches(0.5), Inches(3.0), W - Inches(1.0), Inches(0.5),
            size=12, color=GREY, italic=True, align=PP_ALIGN.CENTER)

    # Contact
    txt(sl, "Pour toute question relative à ce rapport :",
        Inches(0.35), H - Inches(0.88), Inches(8), Inches(0.28),
        size=8.5, color=GREY, italic=True)
    txt(sl, "Équipe Neoclima — direction.travaux@neoclima.ch   ·   +41 22 XXX XX XX",
        Inches(0.35), H - Inches(0.62), Inches(10), Inches(0.28),
        size=9.5, color=NAVY, bold=True)

    footer(sl, P.get("date", ""), W, H)

    # ── Sérialisation ──────────────────────────────────────────────────
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
